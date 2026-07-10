"""PPT/PPTX import: format detection, async dispatch, slide conversion, gating."""

from __future__ import annotations

from pathlib import Path

import pytest
from django.contrib.auth import get_user_model
from django.core.files.base import ContentFile
from django.core.files.uploadedfile import SimpleUploadedFile
from django.urls import reverse
from rest_framework.test import APIClient

from apps.editor import tasks as pptx_tasks
from apps.editor.models import Attachment, SlideImage
from apps.knowledge.models import Document, KnowledgeBase
from apps.knowledge.serializers import detect_doc_format

User = get_user_model()

PPTX_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def _valid_pptx_bytes() -> bytes:
    """Minimal structurally-valid zip (has an EOCD) to pass upload validation.

    Not a real presentation — conversion is monkeypatched in these tests — but a
    loadable zip so the ``_is_valid_zip`` upload guard accepts it.
    """
    import io
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("[Content_Types].xml", "<Types/>")
        z.writestr("ppt/presentation.xml", "<presentation/>")
    return buf.getvalue()


# Corrupt: valid local-file header but the central directory / EOCD is zeroed,
# exactly the "half-downloaded" shape that reaches LibreOffice as an unloadable
# file. ``zipfile.is_zipfile`` returns False for this.
_CORRUPT_PPTX = b"PK\x03\x04" + b"body-bytes" * 50 + b"\x00" * 64


def _png_file(path: Path, color=(30, 120, 200), size=(320, 180)) -> Path:
    from PIL import Image

    Image.new("RGB", size, color).save(path, format="PNG")
    return path


@pytest.fixture
def api_client():
    return APIClient()


@pytest.fixture
def owner():
    return User.objects.create_user("pptowner", "ppt@e.com", "pass", is_staff=True)


@pytest.fixture
def kb(owner):
    return KnowledgeBase.objects.create(
        owner=owner, name="PPT KB", slug="ppt-kb", visibility="public"
    )


# --------------------------------------------------------------------------- #
# format detection
# --------------------------------------------------------------------------- #


@pytest.mark.django_db
def test_detect_doc_format_pptx(owner, kb):
    doc = Document.objects.create(knowledge_base=kb, title="Deck", status="published")
    Attachment.objects.create(
        document=doc,
        uploaded_by=owner,
        file=ContentFile(b"x", name="deck.pptx"),
        original_filename="deck.pptx",
        kind=Attachment.KIND_DOCUMENT,
        mime_type=PPTX_CT,
        size=1,
    )
    assert detect_doc_format(doc) == "pptx"


@pytest.mark.django_db
def test_detect_doc_format_legacy_ppt(owner, kb):
    doc = Document.objects.create(knowledge_base=kb, title="Old", status="published")
    Attachment.objects.create(
        document=doc,
        uploaded_by=owner,
        file=ContentFile(b"x", name="old.ppt"),
        original_filename="old.ppt",
        kind=Attachment.KIND_DOCUMENT,
        mime_type="application/vnd.ms-powerpoint",
        size=1,
    )
    assert detect_doc_format(doc) == "pptx"


# --------------------------------------------------------------------------- #
# import dispatches async conversion
# --------------------------------------------------------------------------- #


@pytest.mark.django_db
def test_import_pptx_dispatches_conversion(api_client, owner, kb, settings, tmp_path, monkeypatch):
    settings.MEDIA_ROOT = str(tmp_path)
    calls = []
    monkeypatch.setattr(
        pptx_tasks.convert_pptx_to_slides,
        "delay",
        lambda doc_id, att_id: calls.append((doc_id, att_id)),
    )
    api_client.force_authenticate(owner)
    f = SimpleUploadedFile("deck.pptx", _valid_pptx_bytes(), content_type=PPTX_CT)
    resp = api_client.post(
        reverse("api_v1:import-file"),
        data={"knowledge_base": kb.id, "file": f},
        format="multipart",
    )
    assert resp.status_code == 201, resp.content
    doc = Document.objects.get(pk=resp.data["id"])
    assert len(calls) == 1 and calls[0][0] == doc.id  # dispatched once for this doc
    assert detect_doc_format(doc) == "pptx"
    # Marked pending so the reader knows conversion is in flight (not failed).
    doc.refresh_from_db()
    assert doc.slide_status == "pending"
    # Body stays empty — it's a view-only binary.
    assert not (doc.published_content or "").strip()


@pytest.mark.django_db
def test_import_corrupt_pptx_rejected(api_client, owner, kb, settings, tmp_path, monkeypatch):
    """B1: a truncated/corrupt pptx (no EOCD) is rejected up front with a clear
    message, instead of being accepted and failing later in async conversion."""
    settings.MEDIA_ROOT = str(tmp_path)
    dispatched = []
    monkeypatch.setattr(
        pptx_tasks.convert_pptx_to_slides, "delay",
        lambda *a: dispatched.append(a),
    )
    api_client.force_authenticate(owner)
    f = SimpleUploadedFile("broken.pptx", _CORRUPT_PPTX, content_type=PPTX_CT)
    resp = api_client.post(
        reverse("api_v1:import-file"),
        data={"knowledge_base": kb.id, "file": f},
        format="multipart",
    )
    assert resp.status_code == 400, resp.content
    assert "损坏" in resp.data["detail"]
    assert not dispatched  # never enqueued a doomed conversion
    assert not Document.objects.filter(title="broken").exists()


@pytest.mark.django_db
def test_import_corrupt_docx_rejected(api_client, owner, kb, settings, tmp_path):
    """The same guard covers .docx (also a zip container)."""
    settings.MEDIA_ROOT = str(tmp_path)
    api_client.force_authenticate(owner)
    f = SimpleUploadedFile(
        "broken.docx", _CORRUPT_PPTX,
        content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    resp = api_client.post(
        reverse("api_v1:import-file"),
        data={"knowledge_base": kb.id, "file": f},
        format="multipart",
    )
    assert resp.status_code == 400, resp.content
    assert "损坏" in resp.data["detail"]


# --------------------------------------------------------------------------- #
# conversion task
# --------------------------------------------------------------------------- #


def _make_doc_with_pptx(owner, kb, tmp_path, settings):
    settings.MEDIA_ROOT = str(tmp_path)
    doc = Document.objects.create(knowledge_base=kb, title="Deck", status="published")
    att = Attachment.objects.create(
        document=doc,
        uploaded_by=owner,
        file=ContentFile(b"PK fake", name="deck.pptx"),
        original_filename="deck.pptx",
        kind=Attachment.KIND_DOCUMENT,
        mime_type=PPTX_CT,
        size=7,
    )
    return doc, att


@pytest.mark.django_db
def test_convert_pptx_creates_ordered_slides(owner, kb, settings, tmp_path, monkeypatch):
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def fake_convert(pptx_path, workdir):
        # Stand in for LibreOffice + pdftoppm: emit 3 ordered PNGs.
        return [
            _png_file(Path(workdir) / f"slide-{i}.png", size=(300 + i, 200))
            for i in range(1, 4)
        ]

    monkeypatch.setattr(pptx_tasks, "_convert", fake_convert)
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 3
    slides = list(SlideImage.objects.filter(document=doc).order_by("index"))
    assert [s.index for s in slides] == [0, 1, 2]
    assert all(s.width > 0 and s.height > 0 for s in slides)
    assert all(s.image for s in slides)
    # Every slide gets a light rail thumbnail; thumb_url points at it (not the raster).
    assert all(s.thumbnail for s in slides)
    assert all(s.thumb_url == s.thumbnail.url for s in slides)
    assert slides[0].source_id == att.id
    doc.refresh_from_db()
    assert doc.slide_status == "done" and doc.slide_error == ""


@pytest.mark.django_db
def test_convert_pptx_idempotent(owner, kb, settings, tmp_path, monkeypatch):
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)
    SlideImage.objects.create(document=doc, source=att, index=0, width=1, height=1,
                              image=ContentFile(b"x", name="s0.png"))
    called = {"n": 0}

    def fake_convert(pptx_path, workdir):
        called["n"] += 1
        return []

    monkeypatch.setattr(pptx_tasks, "_convert", fake_convert)
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 0
    assert called["n"] == 0  # short-circuited before running conversion


@pytest.mark.django_db
def test_convert_pptx_missing_binary_is_soft(owner, kb, settings, tmp_path, monkeypatch):
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def boom(pptx_path, workdir):
        raise FileNotFoundError("soffice")

    monkeypatch.setattr(pptx_tasks, "_convert", boom)
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 0
    assert not SlideImage.objects.filter(document=doc).exists()
    doc.refresh_from_db()
    assert doc.slide_status == "failed"
    assert "组件" in doc.slide_error  # human hint about missing LibreOffice/poppler


@pytest.mark.django_db
def test_convert_pptx_corrupt_sets_failed_reason(owner, kb, settings, tmp_path, monkeypatch):
    """B2: a deck LibreOffice can't load flips slide_status→failed with a reason
    the reader shows (instead of spinning '转换中' forever)."""
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def no_pdf(pptx_path, workdir):
        raise RuntimeError("LibreOffice produced no PDF: Error: source file could not be loaded")

    monkeypatch.setattr(pptx_tasks, "_convert", no_pdf)
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 0
    doc.refresh_from_db()
    assert doc.slide_status == "failed"
    assert "损坏" in doc.slide_error


@pytest.mark.django_db
def test_convert_pptx_no_pages_sets_render_reason(owner, kb, settings, tmp_path, monkeypatch):
    """A raster pass that yields no pages surfaces the specific 光栅化 reason.

    Regression: _failure_reason must track the JPEG output message
    ("pdftoppm produced no JPEG pages"); the old "no PNG" matcher went dead after
    the PNG→JPEG switch and silently degraded this to the generic reason.
    """
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def no_pages(pptx_path, workdir):
        raise RuntimeError("pdftoppm produced no JPEG pages")

    monkeypatch.setattr(pptx_tasks, "_convert", no_pages)
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 0
    doc.refresh_from_db()
    assert doc.slide_status == "failed"
    assert "光栅化" in doc.slide_error


# --------------------------------------------------------------------------- #
# public slides endpoint (friend-gated, ordered)
# --------------------------------------------------------------------------- #


@pytest.mark.django_db
def test_public_slides_endpoint(api_client, owner, kb, settings, tmp_path):
    settings.MEDIA_ROOT = str(tmp_path)
    settings.SITE_REQUIRE_LOGIN = False  # allow anon to reach the gate for this test
    doc = Document.objects.create(
        knowledge_base=kb, title="Deck", status="published", visibility="public",
    )
    att = Attachment.objects.create(
        document=doc, uploaded_by=owner,
        file=ContentFile(b"x", name="deck.pptx"), original_filename="deck.pptx",
        kind=Attachment.KIND_DOCUMENT, mime_type=PPTX_CT, size=1,
    )
    for i in range(2):
        SlideImage.objects.create(
            document=doc, source=att, index=i, width=10 + i, height=20,
            image=ContentFile(b"x", name=f"s{i}.png"),
        )
    doc.slide_status = "done"
    doc.save(update_fields=["slide_status"])
    resp = api_client.get(reverse("api_v1:public-post-slides", args=[doc.id]))
    assert resp.status_code == 200, resp.content
    slides = resp.data["slides"]
    assert [s["index"] for s in slides] == [0, 1]
    assert slides[0]["width"] == 10
    # Legacy rows (no thumbnail) fall back to the full raster url for the rail.
    assert slides[0]["thumb"] == slides[0]["url"]
    # Reader uses these to stop polling / show a real reason.
    assert resp.data["slide_status"] == "done"
    assert resp.data["slide_error"] == ""
    # Notes ship in the slide projection (empty for these note-less legacy rows).
    assert "notes" in slides[0]


# --------------------------------------------------------------------------- #
# speaker notes
# --------------------------------------------------------------------------- #


def _build_pptx_with_notes(path: Path, notes: list[str | None]) -> Path:
    """Write a real .pptx with one slide per entry; a non-None entry gets notes."""
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text in notes:
        slide = prs.slides.add_slide(blank)
        if text is not None:
            slide.notes_slide.notes_text_frame.text = text
    prs.save(str(path))
    return path


def test_extract_pptx_notes_reads_per_slide_in_order(tmp_path):
    p = _build_pptx_with_notes(tmp_path / "deck.pptx", ["first note", None, "third note"])
    assert pptx_tasks.extract_pptx_notes(p) == ["first note", "", "third note"]


def test_extract_pptx_notes_bad_file_is_soft(tmp_path):
    # Unreadable / non-pptx path must yield [] rather than raise (notes are secondary).
    (tmp_path / "junk.pptx").write_bytes(b"not a zip")
    assert pptx_tasks.extract_pptx_notes(tmp_path / "junk.pptx") == []
    assert pptx_tasks.extract_pptx_notes(tmp_path / "missing.pptx") == []


@pytest.mark.django_db
def test_convert_pptx_assigns_notes_by_index(owner, kb, settings, tmp_path, monkeypatch):
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def fake_convert(pptx_path, workdir):
        return [
            _png_file(Path(workdir) / f"slide-{i}.png", size=(300, 200))
            for i in range(1, 4)
        ]

    monkeypatch.setattr(pptx_tasks, "_convert", fake_convert)
    # Notes align to pages by index; middle page has none.
    monkeypatch.setattr(
        pptx_tasks, "extract_pptx_notes", lambda p: ["note A", "", "note C"]
    )
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 3
    slides = list(SlideImage.objects.filter(document=doc).order_by("index"))
    assert [s.notes for s in slides] == ["note A", "", "note C"]
    assert slides[0].as_dict()["notes"] == "note A"


@pytest.mark.django_db
def test_convert_pptx_more_pages_than_notes_leaves_blank(owner, kb, settings, tmp_path, monkeypatch):
    """Hidden-slide drift: fewer notes than rendered pages must not IndexError —
    surplus pages just get empty notes."""
    doc, att = _make_doc_with_pptx(owner, kb, tmp_path, settings)

    def fake_convert(pptx_path, workdir):
        return [_png_file(Path(workdir) / f"slide-{i}.png") for i in range(1, 4)]

    monkeypatch.setattr(pptx_tasks, "_convert", fake_convert)
    monkeypatch.setattr(pptx_tasks, "extract_pptx_notes", lambda p: ["only one"])
    n = pptx_tasks.convert_pptx_to_slides(doc.id, att.id)
    assert n == 3
    slides = list(SlideImage.objects.filter(document=doc).order_by("index"))
    assert [s.notes for s in slides] == ["only one", "", ""]
