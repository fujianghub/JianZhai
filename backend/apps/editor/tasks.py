"""Celery task: convert an uploaded PPT/PPTX into ordered slide images.

Pipeline (see plan F3): LibreOffice ``soffice --headless --convert-to pdf``
turns the presentation into a PDF (layout-faithful, page-ordered), then poppler
``pdftoppm -jpeg`` rasterises each page. Each page becomes a :class:`SlideImage`
row the blog reader renders with a thumbnail rail.

Requires ``libreoffice`` (soffice) + ``poppler-utils`` (pdftoppm) on PATH. When
either binary is missing the task logs a clear hint and leaves slides empty —
the reader keeps showing a "转换中" placeholder rather than crashing.
"""

from __future__ import annotations

import logging
import subprocess
import tempfile
from pathlib import Path

from celery import shared_task
from django.core.files.base import ContentFile
from django.db import transaction

log = logging.getLogger(__name__)

# LibreOffice cold-start + render can be slow for large decks.
_SOFFICE_TIMEOUT = 180
_PDFTOPPM_TIMEOUT = 180
_RASTER_DPI = 150
# Slides are photo-heavy; JPEG cuts a title raster from ~2 MB (PNG) to ~0.27 MB
# with no visible loss, so a 94-slide deck drops from ~24 MB to a few MB.
_JPEG_QUALITY = 82
# Rail thumbnail long-edge in px (displayed at 160px, 2x for retina). ~15-35 KB
# each vs the full raster — the rail was the reader's real weight.
_THUMB_LONG_EDGE = 320
_THUMB_QUALITY = 75


def _run(cmd: list[str], *, timeout: int, cwd: str | None = None) -> str:
    """Run a subprocess, raising on non-zero exit. Returns captured stdout.

    Note: LibreOffice frequently exits 0 while silently refusing to convert a
    file (e.g. "Error: source file could not be loaded" for a corrupt deck), so
    callers must also check the *side effects* — the returned text is what lets
    them surface that hidden message when no output file appears.
    """
    # LibreOffice needs a writable profile dir; point HOME at the temp workspace
    # so it never touches the service user's real home (read-only in Docker).
    env = {"HOME": cwd or tempfile.gettempdir(), "PATH": _path_env()}
    proc = subprocess.run(
        cmd,
        cwd=cwd,
        env=env,
        timeout=timeout,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
    )
    out = (proc.stdout or b"").decode("utf-8", "replace")
    if proc.returncode != 0:
        raise RuntimeError(f"{cmd[0]} exited {proc.returncode}: {out[-1000:]}")
    return out


def _path_env() -> str:
    import os

    return os.environ.get("PATH", "/usr/local/bin:/usr/bin:/bin")


def _convert(pptx_path: Path, workdir: Path) -> list[Path]:
    """Convert a pptx file to a sorted list of per-page JPEG paths."""
    out = _run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(workdir),
            str(pptx_path),
        ],
        timeout=_SOFFICE_TIMEOUT,
        cwd=str(workdir),
    )
    pdfs = list(workdir.glob("*.pdf"))
    if not pdfs:
        # soffice exits 0 even when it can't load the source (corrupt/invalid
        # deck). Its real complaint ("Error: source file could not be loaded")
        # is only in stdout — include it so the failure reason isn't a mystery.
        raise RuntimeError(f"LibreOffice produced no PDF: {out.strip()[-500:]}")
    pdf_path = pdfs[0]

    prefix = workdir / "slide"
    # JPEG (not PNG): slide rasters are photo-heavy, JPEG is ~7x smaller with no
    # visible loss, which is what keeps a large deck light in the reader.
    _run(
        ["pdftoppm", "-jpeg", "-jpegopt", f"quality={_JPEG_QUALITY}",
         "-r", str(_RASTER_DPI), str(pdf_path), str(prefix)],
        timeout=_PDFTOPPM_TIMEOUT,
        cwd=str(workdir),
    )
    # pdftoppm names pages slide-1.jpg, slide-2.jpg … (zero-padded for big decks).
    imgs = sorted(
        workdir.glob("slide-*.jpg"),
        key=lambda p: int(p.stem.rsplit("-", 1)[-1]),
    )
    if not imgs:
        raise RuntimeError("pdftoppm produced no JPEG pages")
    return imgs


def _set_slide_state(document_id: int, status: str, error: str = "") -> None:
    """Persist the doc's slide-conversion state so the reader can stop guessing.

    Best-effort: a status write must never be the thing that crashes the task.
    """
    from apps.knowledge.models import Document

    try:
        Document.objects.filter(pk=document_id).update(
            slide_status=status, slide_error=error[:200]
        )
    except Exception:  # noqa: BLE001
        log.exception("pptx convert: failed to set slide_status for %s", document_id)


# Maps a raw failure to a short, human-facing reason shown in the reader.
def _failure_reason(exc: Exception) -> str:
    msg = str(exc)
    if isinstance(exc, subprocess.TimeoutExpired):
        return "转换超时（文件过大或幻灯过多），请稍后重试"
    if "could not be loaded" in msg or "produced no PDF" in msg:
        return "文件已损坏或无法解析（可能不是有效的 PPT），请重新导出后上传"
    if "pdftoppm produced no" in msg:
        return "幻灯渲染失败（PDF 光栅化未产出页面）"
    return "PPT 转换失败，请重试或联系管理员"


def extract_pptx_notes(pptx_path) -> list[str]:
    """Per-slide speaker notes in presentation order (via python-pptx).

    Best-effort: any failure (python-pptx missing, unreadable pptx) yields ``[]``
    so notes never block the raster conversion — they are secondary content.
    Index ``i`` is meant to align with the i-th rendered PDF page; a deck with
    hidden slides can drift (LibreOffice may skip them in the PDF while
    python-pptx still lists them), so callers map by index and tolerate a length
    mismatch rather than assuming a 1:1 correspondence.
    """
    try:
        from pptx import Presentation
    except ImportError:
        log.warning("python-pptx not installed — skipping PPT notes extraction")
        return []
    try:
        prs = Presentation(str(pptx_path))
    except Exception:  # noqa: BLE001 — a bad deck must not fail the raster path
        log.exception("pptx notes: could not open %s", pptx_path)
        return []
    notes: list[str] = []
    for slide in prs.slides:
        text = ""
        try:
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                if tf is not None:
                    text = (tf.text or "").strip()
        except Exception:  # noqa: BLE001 — skip a single unreadable notes slide
            text = ""
        notes.append(text)
    return notes


@shared_task(name="editor.convert_pptx")
def convert_pptx_to_slides(document_id: int, attachment_id: int) -> int:
    """Render a pptx attachment into ordered SlideImage rows. Returns slide count."""
    from .models import Attachment, SlideImage

    # Idempotency: if slides already exist for this doc, a prior run (or retry)
    # already did the work — don't duplicate.
    if SlideImage.objects.filter(document_id=document_id).exists():
        log.info("pptx %s already has slides — skip", document_id)
        _set_slide_state(document_id, "done")
        return 0

    try:
        att = Attachment.objects.get(pk=attachment_id, document_id=document_id)
    except Attachment.DoesNotExist:
        log.warning("pptx convert: attachment %s not found", attachment_id)
        _set_slide_state(document_id, "failed", "附件丢失，无法转换")
        return 0

    try:
        with tempfile.TemporaryDirectory() as tmp:
            workdir = Path(tmp)
            src_name = Path(att.original_filename or "deck.pptx").name
            pptx_path = workdir / src_name
            with att.file.open("rb") as fh:
                pptx_path.write_bytes(fh.read())

            pages = _convert(pptx_path, workdir)

            # Speaker notes align to rendered pages by index (best-effort — see
            # extract_pptx_notes; a hidden-slide drift just leaves some pages blank).
            notes = extract_pptx_notes(pptx_path)

            import io

            from PIL import Image

            rows = []
            for idx, page in enumerate(pages):
                data = page.read_bytes()
                with Image.open(page) as im:
                    w, h = im.size
                    # Downscale the already-rendered raster into a light rail
                    # thumbnail (cheaper than a second pdftoppm pass over the PDF).
                    thumb = im.convert("RGB")
                    thumb.thumbnail(
                        (_THUMB_LONG_EDGE, _THUMB_LONG_EDGE), Image.LANCZOS
                    )
                    tbuf = io.BytesIO()
                    thumb.save(tbuf, format="JPEG", quality=_THUMB_QUALITY)
                slide = SlideImage(
                    document_id=document_id,
                    source=att,
                    index=idx,
                    width=w,
                    height=h,
                    notes=notes[idx] if idx < len(notes) else "",
                )
                slide.image.save(f"slide-{idx}.jpg", ContentFile(data), save=False)
                slide.thumbnail.save(
                    f"thumb-{idx}.jpg", ContentFile(tbuf.getvalue()), save=False
                )
                rows.append(slide)

            with transaction.atomic():
                # Re-check under nothing fancy; unique_together guards duplicates.
                SlideImage.objects.bulk_create(rows)
            log.info("pptx %s → %d slides", document_id, len(rows))
            _set_slide_state(document_id, "done")
            return len(rows)
    except FileNotFoundError as exc:
        log.error(
            "pptx convert: missing system binary (%s). Install `libreoffice` + "
            "`poppler-utils` and restart the celery worker.",
            exc,
        )
        _set_slide_state(
            document_id, "failed", "服务器缺少 PPT 转换组件（LibreOffice/poppler），请联系管理员"
        )
        return 0
    except Exception as exc:  # noqa: BLE001 — never let a bad deck kill the worker
        log.exception("pptx convert failed for document %s", document_id)
        _set_slide_state(document_id, "failed", _failure_reason(exc))
        return 0


@shared_task
def mirror_document_images(document_id: int, uploaded_by_id: int | None = None) -> dict:
    """Download a markdown document's remote images into local /media storage.

    Runs off the upload request: a Yuque export can carry 40+ images on a CDN
    (cdn.nlark.com) that both anti-hotlinks (browser loads 403 on a foreign
    Referer) and throttles per-IP, so mirroring them synchronously blew past the
    request timeout and left the body full of broken remote URLs. As a task it
    can take its time; the reader shows the remote images (via
    ``referrerpolicy=no-referrer``) until this swaps them for local copies.
    """
    from apps.knowledge.models import Document
    from apps.editor.services.image_mirror import mirror_images_for_document

    doc = Document.objects.filter(pk=document_id).first()
    if doc is None:
        log.info("mirror_document_images: document %s gone, skipping", document_id)
        return {"ok": False, "reason": "document missing"}

    uploaded_by = None
    if uploaded_by_id is not None:
        from django.contrib.auth import get_user_model

        uploaded_by = get_user_model().objects.filter(pk=uploaded_by_id).first()

    result = mirror_images_for_document(doc, uploaded_by=uploaded_by)
    log.info("mirror_document_images %s → %s", document_id, result)
    return result
